# Copyright 2026
# PaperBanana User Manual to PPT Generator

"""
PaperBanana 사용 설명서를 PPT로 변환하는 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os


def create_ppt_from_manual(output_path="PaperBanana_사용설명서.pptx"):
    """PaperBanana 사용 설명서를 PPT로 생성"""

    # 프레젠테이션 생성 (16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 색상 정의
    PRIMARY_COLOR = RGBColor(0, 102, 204)  # 파란색
    SECONDARY_COLOR = RGBColor(255, 193, 7)  # 노란색 (바나나 색상)
    TEXT_COLOR = RGBColor(51, 51, 51)  # 어두운 회색
    BG_COLOR = RGBColor(255, 255, 255)  # 흰색

    def add_title_slide(prs, title, subtitle):
        """제목 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        # 배경
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = PRIMARY_COLOR
        background.line.fill.background()

        # 바나나 이모지
        emoji_box = slide.shapes.add_textbox(
            Inches(0), Inches(2), prs.slide_width, Inches(1)
        )
        tf = emoji_box.text_frame
        tf.text = "🍌"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(72)

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 부제목
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(1)
        )
        tf = subtitle_box.text_frame
        tf.text = subtitle
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)

        return slide

    def add_section_slide(prs, section_title):
        """섹션 구분 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 배경
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = SECONDARY_COLOR
        background.line.fill.background()

        # 섹션 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = section_title
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR

        return slide

    def add_content_slide(prs, title, content_items, has_subtitle=False, subtitle=""):
        """내용 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목 배경
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = PRIMARY_COLOR
        title_bg.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 내용
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5), Inches(12), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        if has_subtitle and subtitle:
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.italic = True
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(12)

        for i, item in enumerate(content_items):
            if i == 0 and not has_subtitle:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)

        return slide

    def add_code_slide(prs, title, code_lines):
        """코드 예제 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 제목 배경
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(45, 45, 45)
        title_bg.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 코드 블록 배경
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3),
            Inches(1.5),
            Inches(12.7),
            Inches(5.6),
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = RGBColor(40, 44, 52)
        code_bg.line.fill.background()

        # 코드 텍스트
        code_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2)
        )
        tf = code_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(code_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line
            p.font.size = Pt(14)
            p.font.name = "Consolas"
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(4)

        return slide

    # ====== 슬라이드 생성 시작 ======

    # 1. 제목 슬라이드
    add_title_slide(
        prs,
        "PaperBanana 사용 설명서",
        "학술 논문 도표 자동 생성 AI 프레임워크\n완벽 가이드",
    )

    # 2. 목차
    add_content_slide(
        prs,
        "📋 목차",
        [
            "PaperBanana 소개 및 주요 특징",
            "프로젝트 구조",
            "빠른 시작 가이드 (설치 및 설정)",
            "사용 방법 - Streamlit 웹 인터페이스",
            "사용 방법 - CLI",
            "실제 사용 예시",
            "고급 설정 및 팁",
            "문제 해결",
            "지원 및 라이선스",
        ],
    )

    # 3. 소개 섹션
    add_section_slide(prs, "🍌 PaperBanana 소개")

    add_content_slide(
        prs,
        "PaperBanana란 무엇인가?",
        [
            "학술 논문의 그림과 도표를 자동으로 생성하는 AI 기반 멀티-에이전트 프레임워크",
            "5개의 전문 에이전트가 협업: Retriever, Planner, Stylist, Visualizer, Critic",
            "논문의 메소드 섹션을 입력받아 출판 품질의 학술 도표 생성",
            "참조 기반 학습으로 유사한 학술 그림들의 스타일과 구조 학습",
            "반복적 개선을 통해 품질 향상",
        ],
        has_subtitle=True,
        subtitle="AI-Powered Academic Illustration Generation",
    )

    add_content_slide(
        prs,
        "주요 특징",
        [
            "🎨 자동 학술 도표 생성 - 논문 메소드를 입력하면 전문적인 도표 자동 생성",
            "🔍 참조 기반 학습 - 유사한 학술 그림들을 참조하여 스타일과 구조 학습",
            "🔄 반복적 개선 - Critic 에이전트가 생성 결과를 평가하고 개선",
            "🎯 두 가지 모드 지원 - 개념 다이어그램(diagram)과 통계 그래프(plot)",
            "🖥️ 직관적 인터페이스 - Streamlit 기반 웹 UI 제공",
            "📊 병렬 처리 - 한 번에 최대 20개 후보 도표 동시 생성",
        ],
    )

    # 4. 구조 섹션
    add_section_slide(prs, "📁 프로젝트 구조")

    add_content_slide(
        prs,
        "주요 디렉토리 및 파일",
        [
            "configs/ - API 키 및 모델 설정",
            "agents/ - AI 에이전트 모듈 (Retriever, Planner, Stylist, Visualizer, Critic, Polish)",
            "data/ - 데이터셋 디렉토리 (PaperBananaBench)",
            "prompts/ - 에이전트 프롬프트",
            "utils/ - 유틸리티 함수",
            "results/ - 생성 결과 저장",
            "main.py - CLI 실행 스크립트",
            "demo.py - Streamlit 웹 UI",
        ],
    )

    # 5. 설치 섹션
    add_section_slide(prs, "🚀 빠른 시작")

    add_code_slide(
        prs,
        "1단계: 환경 설정",
        [
            "# 저장소 클론",
            "git clone https://github.com/dwzhu-pku/PaperBanana.git",
            "cd PaperBanana",
            "",
            "# uv 설치 (아직 설치되지 않은 경우)",
            "curl -LsSf https://astral.sh/uv/install.sh | sh",
            "",
            "# 가상 환경 생성 및 Python 설치",
            "uv venv",
            "uv python install 3.12",
            "",
            "# 의존성 패키지 설치",
            "uv pip install -r requirements.txt",
        ],
    )

    add_code_slide(
        prs,
        "2단계: API 키 설정",
        [
            "# configs/model_config.yaml 파일 수정:",
            "",
            "defaults:",
            '  model_name: "gemini-1.5-pro-preview"',
            '  image_model_name: "gemini-1.5-pro-image-preview"',
            "",
            "api_keys:",
            '  google_api_key: "YOUR_GOOGLE_API_KEY_HERE"',
            '  openai_api_key: "YOUR_OPENAI_API_KEY_HERE"      # 선택사항',
            '  anthropic_api_key: "YOUR_ANTHROPIC_API_KEY_HERE" # 선택사항',
        ],
    )

    add_content_slide(
        prs,
        "API 키 발급 방법",
        [
            "Google Gemini API:",
            "  1. https://makersuite.google.com/app/apikey 접속",
            '  2. "Create API Key" 클릭',
            "  3. 키 복사하여 설정 파일에 붙여넣기",
            "",
            "OpenAI API (선택사항):",
            "  1. https://platform.openai.com/api-keys 접속",
            '  2. "Create new secret key" 클릭',
            "  3. 키 복사하여 설정",
        ],
    )

    # 6. 사용 방법 - Streamlit
    add_section_slide(prs, "💻 사용 방법 - Streamlit")

    add_code_slide(
        prs,
        "Streamlit 웹 인터페이스 실행",
        [
            "# 가상 환경 활성화 후 실행",
            "source .venv/bin/activate  # Linux/Mac",
            ".venv\\Scripts\\activate     # Windows",
            "",
            "# Streamlit 실행",
            "streamlit run demo.py",
            "",
            "# 브라우저에서 http://localhost:8501 로 접속",
        ],
    )

    add_content_slide(
        prs,
        "탭 1: Generate Candidates (후보 생성)",
        [
            "Pipeline Mode 선택:",
            "  • demo_planner_critic: Planner → Visualizer → Critic → Visualizer",
            "  • demo_full: Retriever → Planner → Stylist → Visualizer → Critic → Visualizer",
            "",
            "설정 옵션:",
            "  • Number of Candidates: 생성할 후보 수 (1-20)",
            "  • Aspect Ratio: 이미지 비율 (21:9, 16:9, 3:2)",
            "  • Max Critic Rounds: 개선 반복 횟수 (1-5)",
        ],
    )

    add_content_slide(
        prs,
        "입력 및 결과",
        [
            "입력:",
            "  • Method Section Content: 논문의 메소드 섹션 (Markdown 권장)",
            "  • Figure Caption: 생성할 그림의 캡션/설명",
            "",
            "결과:",
            "  • 여러 후보 이미지를 그리드로 표시",
            "  • 각 후보의 진화 타임라인 보기",
            "  • 개별 이미지 또는 ZIP으로 다운로드",
        ],
    )

    # 7. CLI
    add_section_slide(prs, "⌨️ 사용 방법 - CLI")

    add_code_slide(
        prs,
        "명령줄 인터페이스",
        [
            "# 기본 사용법 (기본 설정)",
            "python main.py",
            "",
            "# 고급 사용법 (사용자 설정)",
            "python main.py",
            '  --dataset_name "PaperBananaBench"',
            '  --task_name "diagram"',
            '  --split_name "test"',
            '  --exp_mode "dev_full"',
            '  --retrieval_setting "auto"',
            "  --max_critic_rounds 3",
        ],
    )

    add_content_slide(
        prs,
        "CLI 옵션",
        [
            "--dataset_name: 사용할 데이터셋 이름 (기본값: PaperBananaBench)",
            "--task_name: 작업 유형 (diagram 또는 plot)",
            "--split_name: 데이터셋 분할 (test/train/val)",
            "--exp_mode: 실험 모드 (vanilla, dev_planner, dev_full 등)",
            "--retrieval_setting: 검색 설정 (auto/manual/random/none)",
            "--max_critic_rounds: 최대 Critic 반복 횟수",
            "--model_name: 사용할 모델 이름",
        ],
    )

    # 8. 실험 모드
    add_content_slide(
        prs,
        "실험 모드 종류",
        [
            "vanilla: Visualizer만 사용 (기본 생성)",
            "dev_planner: Planner → Visualizer",
            "dev_planner_stylist: Planner → Stylist → Visualizer",
            "dev_planner_critic: Planner → Visualizer → Critic (다중 라운드)",
            "dev_full: 전체 파이프라인 (모든 에이전트)",
            "demo_planner_critic: 데모용 (평가 없음)",
            "demo_full: 데모용 전체 파이프라인 (평가 없음)",
        ],
    )

    # 9. 사용 예시
    add_section_slide(prs, "🎯 실제 사용 예시")

    add_content_slide(
        prs,
        "예시: Transformer 아키텍처 다이어그램",
        [
            "Method Section: Transformer 아키텍처 설명",
            "  • Encoder/Decoder 구조",
            "  • Multi-head Attention 메커니즘",
            "  • Residual connections",
            "",
            "Caption: The Transformer architecture 설명",
            "",
            "설정:",
            "  • Pipeline Mode: demo_full",
            "  • Aspect Ratio: 16:9",
            "  • Candidates: 10",
            "  • Critic Rounds: 3",
        ],
    )

    # 10. 팁
    add_section_slide(prs, "💡 고급 설정 및 팁")

    add_content_slide(
        prs,
        "최상의 결과를 위한 팁",
        [
            "메소드 섹션 작성 가이드:",
            "  ✅ Markdown 형식 사용 (제목, 목록 등 구조화)",
            "  ✅ 수식은 LaTeX 형식으로 작성 ($E=mc^2$)",
            "  ✅ 컴포넌트와 흐름을 명확히 설명",
            "  ❌ 너무 긴 배경 설명은 피하기",
            "  ❌ 모호한 표현 피하기",
            "",
            "캡션 작성 가이드:",
            "  • 생성하려는 그림의 핵심 내용 요약",
        ],
    )

    add_content_slide(
        prs,
        "파이프라인 모드 선택 가이드",
        [
            "빠른 결과 필요: demo_planner_critic",
            "  → 빠륵지만 품질 좋음",
            "",
            "최고 품질 필요: demo_full",
            "  → Stylist로 미학적 개선",
            "",
            "단순 도표: vanilla",
            "  → 가장 빠름",
            "",
            "연구/개발: dev_* 모드",
            "  → 평가 메트릭 포함",
        ],
    )

    # 11. 문제 해결
    add_section_slide(prs, "🐛 문제 해결")

    add_content_slide(
        prs,
        "자주 발생하는 문제",
        [
            '"No module named streamlit" 오류:',
            "  → 가상 환경 활성화 확인",
            "  → uv pip install -r requirements.txt 재실행",
            "",
            "API 키 오류:",
            "  → configs/model_config.yaml 파일 경로 확인",
            "  → 환경 변수로 설정하는 대안 사용 가능",
            "",
            "이미지 생성 실패:",
            "  → API 횟수 제한 확인",
            "  → 입력 텍스트 길이 확인",
            "  → 모델 이름 확인",
        ],
    )

    # 12. 마무리
    add_section_slide(prs, "📞 지원 및 라이선스")

    add_content_slide(
        prs,
        "리소스 및 지원",
        [
            "GitHub Issues:",
            "  https://github.com/dwzhu-pku/PaperBanana/issues",
            "",
            "HuggingFace Dataset:",
            "  https://huggingface.co/datasets/dwzhu/PaperBananaBench",
            "",
            "Paper:",
            "  https://huggingface.co/papers/2601.23265",
            "",
            "라이선스: Apache-2.0",
        ],
    )

    # 13. 마지막 슬라이드
    add_title_slide(
        prs,
        "감사합니다!",
        "PaperBanana로 학술 도표 제작을 자동화하세요\n\n🍌 Happy Research! 🍌",
    )

    # PPT 저장
    prs.save(output_path)
    print("[SUCCESS] PPT created successfully!")
    print(f"[INFO] Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    output_file = create_ppt_from_manual()
    print(f"\n[USAGE] Open {output_file} in PowerPoint or Google Slides.")
    print(
        f"   Microsoft PowerPoint 또는 Google Slides에서 {output_file} 파일을 열어주세요."
    )
